from pptx import Presentation
from pptx.util import Pt, Emu
import re

prs = Presentation(r'C:\Work\SR Tech News Daily_0310_수정.pptx')

print("=== SLIDE SIZE ===")
w, h = prs.slide_width, prs.slide_height
print(f"Width: {w/914400:.2f}in ({w/12700:.0f}pt)")
print(f"Height: {h/914400:.2f}in ({h/12700:.0f}pt)")
print(f"Pixel @96dpi: {w/914400*96:.0f} x {h/914400*96:.0f}")
print()

# Analyze slides 2-5
for si in range(1, min(8, len(prs.slides))):
    slide = prs.slides[si]
    print(f"=== SLIDE {si+1} ===")
    
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table = shape.table
            print(f"  TABLE: {len(table.rows)}x{len(table.columns)}")
            print(f"  pos: x={shape.left/914400:.2f} y={shape.top/914400:.2f} w={shape.width/914400:.2f} h={shape.height/914400:.2f}")
            
            for ci, col in enumerate(table.columns):
                print(f"    col[{ci}] w={col.width/914400:.2f}in")
            
            for ri, row in enumerate(table.rows):
                print(f"    row[{ri}] h={row.height/914400:.2f}in")
                for ci, cell in enumerate(row.cells):
                    txt = cell.text[:80].replace('\n', ' | ') if cell.text else ''
                    print(f"      [{ri},{ci}]: {txt}")
                    for pi, para in enumerate(cell.text_frame.paragraphs[:3]):
                        for run in para.runs[:1]:
                            f = run.font
                            try:
                                clr = str(f.color.rgb)
                            except:
                                clr = 'theme'
                            print(f"        font={f.name} sz={f.size} color={clr} bold={f.bold} ul={f.underline}")
            
            # Find colors in XML
            xml = shape._element.xml[:5000]
            colors = re.findall(r'val="([0-9A-Fa-f]{6})"', xml)
            print(f"  Colors in table: {list(set(colors))}")
        
        elif shape.shape_type == 13:  # PICTURE
            print(f"  IMAGE: x={shape.left/914400:.2f} y={shape.top/914400:.2f} w={shape.width/914400:.2f} h={shape.height/914400:.2f}")
        
        elif shape.has_text_frame:
            txt = shape.text[:100].replace('\n', ' | ')
            print(f"  TEXT: '{txt}'")
            print(f"    pos: x={shape.left/914400:.2f} y={shape.top/914400:.2f}")
            for run in shape.text_frame.paragraphs[0].runs[:1]:
                f = run.font
                try:
                    clr = str(f.color.rgb)
                except:
                    clr = 'theme'
                print(f"    font={f.name} sz={f.size} color={clr} bold={f.bold}")
    print()
